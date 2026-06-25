#!/usr/bin/env python3
"""Generate PowerPoint presentation for the QGIS + AI Site Selection demo."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ─── Color palette ───────────────────────────────────────────────
DARK    = RGBColor(0x1A, 0x1A, 0x2E)  # deep navy
MEDIUM  = RGBColor(0x16, 0x21, 0x3E)
BLUE    = RGBColor(0x0F, 0x34, 0x60)
TEAL    = RGBColor(0x4F, 0xC3, 0xF7)
GREEN   = RGBColor(0x31, 0xA3, 0x54)
RED     = RGBColor(0xE7, 0x4C, 0x3C)
ORANGE  = RGBColor(0xF3, 0x9C, 0x12)
PURPLE  = RGBColor(0x8E, 0x44, 0xAD)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT   = RGBColor(0xF0, 0xF2, 0xF5)
GRAY    = RGBColor(0x66, 0x66, 0x66)
DARKTEXT = RGBColor(0x1A, 0x1A, 0x2E)

current_dir = os.path.dirname(os.path.abspath(__file__))

def add_bg(slide, color=DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if alpha is not None:
        shape.fill.fore_color.brightness = alpha
    return shape

def add_accent_bar(slide, left=Inches(0), top=Inches(0), width=Inches(0.08), height=Inches(7.5), color=TEAL):
    return add_shape(slide, left, top, width, height, color)

def add_text_box(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox

def add_bullet_slide(slide, left, top, width, height, items, font_size=20, color=WHITE, spacing=Pt(8), font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = spacing
        p.level = 0
    return txBox

def add_icon_bullets(slide, left, top, width, height, items, font_size=20, color=WHITE):
    """items = [(emoji, text), ...]"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (emoji, text) in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = f"{emoji}  "
        run.font.size = Pt(font_size + 4)
        run2 = p.add_run()
        run2.text = text
        run2.font.size = Pt(font_size)
        run2.font.color.rgb = color
        run2.font.name = "Calibri"
        p.space_after = Pt(12)
    return txBox

def hex_to_rgb(hex_str):
    h = hex_str.lstrip('#')
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

def add_footer(slide, text="QGIS + AI Automation Workshop · github.com/your-repo"):
    add_text_box(slide, Inches(0.8), Inches(6.8), Inches(8), Inches(0.5),
                 text, font_size=10, color=GRAY)

def slide_number(slide, num, total=15):
    add_text_box(slide, Inches(12), Inches(6.9), Inches(1), Inches(0.5),
                 f"{num}/{total}", font_size=10, color=GRAY, align=PP_ALIGN.RIGHT)

def new_slide(num):
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)
    add_bg(slide, DARK)
    add_accent_bar(slide)
    slide_number(slide, num)
    return slide

# ═══════════════════════════════════════════════════════════════════
# SLIDE 1: TITLE
# ═══════════════════════════════════════════════════════════════════
slide = new_slide(1)
add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.5),
             "QGIS + AI", font_size=60, color=WHITE, bold=True)
add_text_box(slide, Inches(1), Inches(3), Inches(11), Inches(1),
             "Automate Your GIS Workflow", font_size=36, color=TEAL, bold=False)
add_shape(slide, Inches(1), Inches(4.2), Inches(3), Inches(0.06), TEAL)
add_text_box(slide, Inches(1), Inches(4.6), Inches(11), Inches(1),
             "Site Selection Analysis with PyQGIS & Python",
             font_size=22, color=GRAY)
add_text_box(slide, Inches(1), Inches(5.6), Inches(11), Inches(1),
             "From 4,708 buildings to 9 suitable sites — in one command",
             font_size=18, color=WHITE, bold=False)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 2: THE PROBLEM
# ═══════════════════════════════════════════════════════════════════
slide = new_slide(2)
add_text_box(slide, Inches(1), Inches(0.5), Inches(11), Inches(1),
             "❌  The Problem", font_size=44, color=WHITE, bold=True)
add_shape(slide, Inches(1), Inches(1.3), Inches(2), Inches(0.06), TEAL)

add_icon_bullets(slide, Inches(1), Inches(2), Inches(5.5), Inches(5), [
    ("🕐", "Manual site selection takes 4+ hours"),
    ("🔄", "Buffer → Clip → Select → Export → Repeat"),
    ("❌", "Error-prone — miss a step, wrong answer"),
    ("📋", "Not reproducible — can you redo it next month?"),
    ("💰", "Expensive — billable hours for repetitive tasks"),
    ("🧠", "Your skill is analysis, not clicking buttons"),
])

# Right side: criteria box
add_shape(slide, Inches(7.2), Inches(2), Inches(5.2), Inches(4.5), MEDIUM)
add_text_box(slide, Inches(7.5), Inches(2.2), Inches(4.7), Inches(0.5),
             "🎯  The Task", font_size=24, color=TEAL, bold=True)
add_icon_bullets(slide, Inches(7.5), Inches(3), Inches(4.7), Inches(3.5), [
    ("1.", "Area between 800–1,200 m²"),
    ("2.", "Within 1 km of a school"),
    ("3.", "NOT within 100 m of a highway"),
    ("4.", "NOT within 300 m of a river"),
    ("5.", "Inside a residential urban zone"),
    ("6.", "Calculate population estimate"),
], font_size=18, color=WHITE)

add_footer(slide)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 3: THE SOLUTION
# ═══════════════════════════════════════════════════════════════════
slide = new_slide(3)
add_text_box(slide, Inches(1), Inches(0.5), Inches(11), Inches(1),
             "✅  The Solution", font_size=44, color=WHITE, bold=True)
add_shape(slide, Inches(1), Inches(1.3), Inches(2), Inches(0.06), GREEN)

add_icon_bullets(slide, Inches(1), Inches(2), Inches(5.5), Inches(5), [
    ("⚡", "QGIS Headless Mode + PyQGIS = Automation"),
    ("🔧", "Same algorithms you know — just scripted"),
    ("📦", "800+ processing algorithms available"),
    ("🔁", "Repeatable · Auditable · Shareable"),
    ("🐍", "One Python script replaces hours of clicking"),
    ("🌍", "Publish to web dashboard automatically"),
])

# Right side: code snippet
add_shape(slide, Inches(7.2), Inches(2), Inches(5.2), Inches(4.5), RGBColor(0x0D, 0x0D, 0x1A))
add_text_box(slide, Inches(7.5), Inches(2.2), Inches(4.7), Inches(0.5),
             "🐍  One Command", font_size=24, color=GREEN, bold=True)
code_text = """  python3 demo.py

  ┌──────────────────────────┐
  │  4,708 buildings         │
  │  → 28  (area filter)     │
  │  → 22  (school zone)     │
  │  → 21  (no highway)      │
  │  →  9  (no river)        │
  │  →  9  (urban zone) ✅   │
  └──────────────────────────┘"""
add_text_box(slide, Inches(7.5), Inches(3), Inches(4.7), Inches(3.5),
             code_text, font_size=15, color=GREEN, font_name="Consolas")

add_footer(slide)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 4: THE DATASET
# ═══════════════════════════════════════════════════════════════════
slide = new_slide(4)
add_text_box(slide, Inches(1), Inches(0.5), Inches(11), Inches(1),
             "🗂️  The Dataset", font_size=44, color=WHITE, bold=True)
add_shape(slide, Inches(1), Inches(1.3), Inches(2), Inches(0.06), BLUE)

# Table as individual cards
layers = [
    ("🏘️  Buildings", "4,708", "Polygon", hex_to_rgb("#e74c3c")),
    ("🏫  Schools", "4", "Polygon", hex_to_rgb("#2980b9")),
    ("🛣️  Roads", "717", "Line", hex_to_rgb("#f39c12")),
    ("🌊  Rivers", "19", "Line", hex_to_rgb("#16a085")),
    ("🏙️  Urban Zones", "6", "Polygon", hex_to_rgb("#8e44ad")),
    ("💧  Water", "37", "Polygon", hex_to_rgb("#3498db")),
    ("🍽️  Restaurants", "11", "Polygon", hex_to_rgb("#e84393")),
]

for i, (name, count, geom, color) in enumerate(layers):
    row = i // 4
    col = i % 4
    x = Inches(1) + col * Inches(3.0)
    y = Inches(2.2) + row * Inches(2.4)
    add_shape(slide, x, y, Inches(2.7), Inches(2), RGBColor(0x16, 0x21, 0x3E))
    # Color indicator
    add_shape(slide, x, y, Inches(0.08), Inches(2), color)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.3), Inches(2.2), Inches(0.5),
                 name, font_size=18, color=WHITE, bold=True)
    add_text_box(slide, x + Inches(0.3), y + Inches(1.0), Inches(2.2), Inches(0.5),
                 f"{count} features", font_size=28, color=color, bold=True)
    add_text_box(slide, x + Inches(0.3), y + Inches(1.5), Inches(2.2), Inches(0.4),
                 f"({geom})", font_size=14, color=GRAY)

add_text_box(slide, Inches(1), Inches(6.5), Inches(11), Inches(0.5),
             "📍 Cape Agulhas, South Africa  ·  CRS: EPSG:4326 / EPSG:32733  ·  Dataset: QGIS Training Manual · Data: © OpenStreetMap (ODbL)",
             font_size=14, color=GRAY)

add_footer(slide)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 5: THE CRITERIA
# ═══════════════════════════════════════════════════════════════════
slide = new_slide(5)
add_text_box(slide, Inches(1), Inches(0.5), Inches(11), Inches(1),
             "📋  The Selection Criteria", font_size=44, color=WHITE, bold=True)
add_shape(slide, Inches(1), Inches(1.3), Inches(2), Inches(0.06), TEAL)

criteria = [
    ("1", "Area 800 – 1,200 m²", "native:extractbyexpression", "Filter buildings by size", GREEN),
    ("2", "Within 1,000 m of school", "native:buffer → extractbylocation", "Buffer schools, select intersect", BLUE),
    ("3", "NOT within 100 m of highway", "native:buffer → extractbylocation", "Buffer highways, select disjoint", RED),
    ("4", "NOT within 300 m of river", "native:buffer → extractbylocation", "Buffer rivers, select disjoint", ORANGE),
    ("5", "Inside urban (residential) zone", "native:extractbylocation", "Select intersect with urban", PURPLE),
    ("6", "Calculate population estimate", "native:fieldcalculator", "Area / 30 = people per building", TEAL),
]

for i, (num, criterion, algorithm, description, color) in enumerate(criteria):
    y = Inches(2) + i * Inches(0.85)
    # Number circle
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1), y + Inches(0.05), Inches(0.5), Inches(0.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.paragraphs[0].text = num
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    add_text_box(slide, Inches(1.8), y, Inches(4), Inches(0.4),
                 criterion, font_size=20, color=WHITE, bold=True)
    add_text_box(slide, Inches(1.8), y + Inches(0.35), Inches(5), Inches(0.3),
                 f"Algorithm: {algorithm}", font_size=13, color=GRAY, font_name="Consolas")
    add_text_box(slide, Inches(7.5), y + Inches(0.05), Inches(5), Inches(0.5),
                 description, font_size=16, color=LIGHT)

add_footer(slide)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 6: THE PIPELINE (Flow diagram)
# ═══════════════════════════════════════════════════════════════════
slide = new_slide(6)
add_text_box(slide, Inches(1), Inches(0.3), Inches(11), Inches(0.8),
             "⚙️  The Processing Pipeline", font_size=40, color=WHITE, bold=True)
add_shape(slide, Inches(1), Inches(1.0), Inches(2), Inches(0.06), TEAL)

# Flow boxes
steps = [
    ("Load GeoPackage\nLayers", BLUE),
    ("Reproject\nto UTM 33S", BLUE),
    ("Filter by\nArea 800-1200m²", GREEN),
    ("Buffer Schools\n1,000m → Select", TEAL),
    ("Buffer Highways\n100m → Exclude", RED),
    ("Buffer Rivers\n300m → Exclude", ORANGE),
    ("Clip to\nUrban Zones", PURPLE),
    ("Calculate\nPopulation", GREEN),
    ("Export\nDashboard", GREEN),
]

box_w = Inches(1.2)
box_h = Inches(1.2)
gap = Inches(0.12)
start_x = Inches(0.4)
start_y = Inches(2.5)

for i, (label, color) in enumerate(steps):
    x = start_x + i * (box_w + gap)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, start_y, box_w, box_h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = label
    tf.paragraphs[0].font.size = Pt(10)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Arrow
    if i < len(steps) - 1:
        arrow_x = x + box_w
        arr = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, arrow_x, start_y + Inches(0.5), Inches(0.12), Inches(0.2))
        arr.fill.solid()
        arr.fill.fore_color.rgb = TEAL
        arr.line.fill.background()

# Stats at bottom
add_text_box(slide, Inches(0.5), Inches(4.3), Inches(12), Inches(0.4),
             "PyQGIS Processing Algorithms (native:*) — same tools as the QGIS toolbox, automated via Python",
             font_size=14, color=GRAY, align=PP_ALIGN.CENTER)

# Before/After boxes
add_shape(slide, Inches(1.5), Inches(5.2), Inches(4.5), Inches(1.5), MEDIUM)
add_text_box(slide, Inches(1.8), Inches(5.4), Inches(4), Inches(0.4),
             "📥  Input: 4,708 buildings", font_size=20, color=WHITE, bold=True)
add_text_box(slide, Inches(1.8), Inches(5.9), Inches(4), Inches(0.5),
             "7 raw GeoPackage layers\n+ satellite imagery", font_size=14, color=GRAY)

add_shape(slide, Inches(7.3), Inches(5.2), Inches(4.5), Inches(1.5), GREEN)
add_text_box(slide, Inches(7.6), Inches(5.4), Inches(4), Inches(0.4),
             "📤  Output: 9 suitable sites", font_size=20, color=WHITE, bold=True)
add_text_box(slide, Inches(7.6), Inches(5.9), Inches(4), Inches(0.5),
             "Interactive dashboard · GeoJSON\nPDF report · Population estimate", font_size=14, color=WHITE)

add_footer(slide)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 7: THE CASCADE (WOW SLIDE)
# ═══════════════════════════════════════════════════════════════════
slide = new_slide(7)
add_text_box(slide, Inches(1), Inches(0.3), Inches(11), Inches(0.8),
             "📊  The Filter Cascade", font_size=40, color=WHITE, bold=True)
add_shape(slide, Inches(1), Inches(1.0), Inches(2), Inches(0.06), TEAL)

cascade_data = [
    ("Total Buildings", 4708, DARK, 1.0),
    ("Area 800-1200 m²", 28, GREEN, 0.006),
    ("Within 1 km of School", 22, TEAL, 0.005),
    ("Not Near Highway", 21, ORANGE, 0.004),
    ("Not Near River", 9, RED, 0.002),
    ("In Urban Zone ✅", 9, GREEN, 0.002),
]

max_bar_width = Inches(9)
for i, (label, count, color, fraction) in enumerate(cascade_data):
    y = Inches(1.6) + i * Inches(0.85)
    bar_width = Inches(max(fraction * 9, 0.5))

    # Label
    add_text_box(slide, Inches(1), y + Inches(0.05), Inches(4), Inches(0.4),
                 label, font_size=18, color=WHITE, bold=(i == 0 or i == len(cascade_data)-1))

    # Count
    count_color = GREEN if i == len(cascade_data) - 1 else WHITE
    add_text_box(slide, Inches(11), y + Inches(0.05), Inches(1.5), Inches(0.4),
                 str(count), font_size=22, color=count_color, bold=True, align=PP_ALIGN.RIGHT)

    # Bar
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(5), y + Inches(0.1), bar_width, Inches(0.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

add_footer(slide)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 8: THE DASHBOARD
# ═══════════════════════════════════════════════════════════════════
slide = new_slide(8)
add_text_box(slide, Inches(1), Inches(0.3), Inches(11), Inches(0.8),
             "🖥️  The Showcase Dashboard", font_size=40, color=WHITE, bold=True)
add_shape(slide, Inches(1), Inches(1.0), Inches(2), Inches(0.06), TEAL)

features = [
    ("📊  Stats Cards", "Filtering cascade at a glance — 4,708 → 9"),
    ("🗺️  Interactive Map", "Green = suitable sites · Red = highway exclusion · Blue = school zone"),
    ("📉  Filter Funnel", "Visual bar chart of each selection step"),
    ("📋  Data Table", "Searchable, sortable, paginated site register"),
    ("📥  Export CSV", "One-click download of the sites register"),
    ("🌐  Export GeoJSON", "One-click download for QGIS/PostGIS"),
    ("📄  Print / PDF", "Ctrl+P → Save as PDF — professional report ready"),
]

for i, (title, desc) in enumerate(features):
    y = Inches(1.6) + i * Inches(0.75)
    add_shape(slide, Inches(1), y, Inches(0.08), Inches(0.6), TEAL)
    add_text_box(slide, Inches(1.3), y, Inches(3.5), Inches(0.4),
                 title, font_size=20, color=TEAL, bold=True)
    add_text_box(slide, Inches(5), y + Inches(0.02), Inches(7), Inches(0.4),
                 desc, font_size=16, color=WHITE)

add_text_box(slide, Inches(1), Inches(6.5), Inches(11), Inches(0.5),
             "💡  Open phase3_dashboard.html in any browser — no installation needed",
             font_size=14, color=GRAY)

add_footer(slide)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 9: LIVE DEMO
# ═══════════════════════════════════════════════════════════════════
slide = new_slide(9)
add_text_box(slide, Inches(1), Inches(0.5), Inches(11), Inches(1),
             "🎬  Live Demo", font_size=44, color=WHITE, bold=True)
add_shape(slide, Inches(1), Inches(1.3), Inches(2), Inches(0.06), TEAL)

add_text_box(slide, Inches(1), Inches(2.5), Inches(11), Inches(1),
             "python3 demo.py", font_size=36, color=GREEN, bold=True, font_name="Consolas",
             align=PP_ALIGN.CENTER)

add_shape(slide, Inches(4), Inches(3.8), Inches(5.3), Inches(0.06), TEAL)

add_icon_bullets(slide, Inches(1.5), Inches(4.2), Inches(10), Inches(3), [
    ("1️⃣", "Open terminal → type python3 demo.py"),
    ("2️⃣", "Press Enter → watch all 7 layers load"),
    ("3️⃣", "Watch the cascade: 4,708 → 9"),
    ("4️⃣", "Dashboard opens automatically in browser"),
    ("5️⃣", "Explore data table, export options, PDF report"),
])

add_footer(slide)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 10: HOW IT WORKS (CODE)
# ═══════════════════════════════════════════════════════════════════
slide = new_slide(10)
add_text_box(slide, Inches(1), Inches(0.3), Inches(11), Inches(0.8),
             "🐍  How It Works — PyQGIS Code", font_size=40, color=WHITE, bold=True)
add_shape(slide, Inches(1), Inches(1.0), Inches(2), Inches(0.06), TEAL)

add_text_box(slide, Inches(1), Inches(1.4), Inches(11), Inches(0.4),
             "Same QGIS toolbox tools — automated in 10 lines of Python",
             font_size=18, color=GRAY)

code = """  # 1. Reproject to UTM 33S for meter-based analysis
  buildings_utm = processing.run("native:reprojectlayer", {
      'INPUT': buildings, 'TARGET_CRS': crs_utm
  })['OUTPUT']

  # 2. Filter buildings by area
  suitable = processing.run("native:extractbyexpression", {
      'INPUT': buildings_utm,
      'EXPRESSION': '"area_sqm" >= 800 AND "area_sqm" <= 1200'
  })['OUTPUT']

  # 3. Buffer schools 1000m, select within
  school_zone = processing.run("native:buffer", {
      'INPUT': schools, 'DISTANCE': 1000
  })['OUTPUT']

  result = processing.run("native:extractbylocation", {
      'INPUT': suitable, 'PREDICATE': [0], 'INTERSECT': school_zone
  })['OUTPUT']"""

add_shape(slide, Inches(1), Inches(2.2), Inches(11.3), Inches(4.5), RGBColor(0x0D, 0x0D, 0x1A))
add_text_box(slide, Inches(1.3), Inches(2.5), Inches(10.7), Inches(4),
             code, font_size=14, color=GREEN, font_name="Consolas")

add_footer(slide)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 11: PRODUCTION ARCHITECTURE
# ═══════════════════════════════════════════════════════════════════
slide = new_slide(11)
add_text_box(slide, Inches(1), Inches(0.3), Inches(11), Inches(0.8),
             "🏗️  Production Architecture", font_size=40, color=WHITE, bold=True)
add_shape(slide, Inches(1), Inches(1.0), Inches(2), Inches(0.06), TEAL)

# Architecture boxes with arrows
arch = [
    ("PyQGIS\nScript", BLUE),
    ("PostgreSQL\n+ PostGIS", MEDIUM),
    ("GeoServer\nWMS / WFS", PURPLE),
    ("Leaflet\nWeb Map", GREEN),
]

for i, (label, color) in enumerate(arch):
    x = Inches(0.8) + i * Inches(3.1)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.5), Inches(2.5), Inches(1.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = label
    tf.paragraphs[0].font.size = Pt(20)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    if i < len(arch) - 1:
        arr = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                      x + Inches(2.5), Inches(3.0), Inches(0.6), Inches(0.3))
        arr.fill.solid()
        arr.fill.fore_color.rgb = TEAL
        arr.line.fill.background()

# Benefits
add_shape(slide, Inches(1), Inches(4.8), Inches(11.3), Inches(2.2), MEDIUM)
add_text_box(slide, Inches(1.3), Inches(5.0), Inches(10.7), Inches(0.4),
             "✨  Benefits", font_size=22, color=TEAL, bold=True)
add_icon_bullets(slide, Inches(1.3), Inches(5.5), Inches(10.7), Inches(1.5), [
    ("🔄", "Scheduled re-analysis — cron / GitHub Actions"),
    ("🌐", "Web services — REST API, tile caching"),
    ("👥", "User permissions — role-based access"),
    ("📱", "Access from any device — desktop, tablet, phone"),
], font_size=16)

add_footer(slide)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 12: USE CASES
# ═══════════════════════════════════════════════════════════════════
slide = new_slide(12)
add_text_box(slide, Inches(1), Inches(0.3), Inches(11), Inches(0.8),
             "🌍  Use Cases", font_size=40, color=WHITE, bold=True)
add_shape(slide, Inches(1), Inches(1.0), Inches(2), Inches(0.06), TEAL)

usecases = [
    ("🏪", "Retail Site Selection", "Find optimal store locations based on demographics, competitor proximity, traffic patterns"),
    ("🏥", "Healthcare Facility Placement", "Identify underserved areas for clinics, hospitals, and emergency services"),
    ("🌳", "Conservation Reserve Design", "Design protected areas using ecological criteria, connectivity, and land use"),
    ("🚒", "Emergency Response Planning", "Locate fire stations, evacuation centers, and supply depots for disaster readiness"),
    ("🏘️", "Urban Development Zoning", "Evaluate parcels for rezoning, density changes, and mixed-use development"),
    ("📡", "Telecom Tower Placement", "Optimize coverage by analyzing terrain, population density, and existing infrastructure"),
]

for i, (emoji, title, desc) in enumerate(usecases):
    row = i // 2
    col = i % 2
    x = Inches(0.8) + col * Inches(6.2)
    y = Inches(1.6) + row * Inches(1.8)
    add_shape(slide, x, y, Inches(5.8), Inches(1.5), MEDIUM)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.2), Inches(5.2), Inches(0.5),
                 f"{emoji}  {title}", font_size=22, color=TEAL, bold=True)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.8), Inches(5.2), Inches(0.6),
                 desc, font_size=14, color=WHITE)

add_footer(slide)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 13: KEY TAKEAWAYS
# ═══════════════════════════════════════════════════════════════════
slide = new_slide(13)
add_text_box(slide, Inches(1), Inches(0.3), Inches(11), Inches(0.8),
             "💡  Key Takeaways", font_size=40, color=WHITE, bold=True)
add_shape(slide, Inches(1), Inches(1.0), Inches(2), Inches(0.06), TEAL)

takeaways = [
    ("⚡", "QGIS + Python replaces hours of clicking — one command, 25 seconds"),
    ("🔧", "PyQGIS gives you access to 800+ processing algorithms"),
    ("♻️", "Same workflow for ANY site selection problem — retail, conservation, emergency"),
    ("📈", "From local script → PostgreSQL → GeoServer → public web map"),
    ("💰", "100% open source — QGIS, Python, Folium, PostgreSQL — no license fees"),
    ("🌍", "Perfect for low-resource regions where ArcGIS is unaffordable"),
]

for i, (emoji, text) in enumerate(takeaways):
    y = Inches(1.6) + i * Inches(0.85)
    add_shape(slide, Inches(1), y, Inches(0.08), Inches(0.6), TEAL)
    add_text_box(slide, Inches(1.3), y + Inches(0.05), Inches(11), Inches(0.5),
                 f"{emoji}  {text}", font_size=18, color=WHITE)

add_footer(slide)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 14: RESOURCES & NEXT STEPS
# ═══════════════════════════════════════════════════════════════════
slide = new_slide(14)
add_text_box(slide, Inches(1), Inches(0.3), Inches(11), Inches(0.8),
             "📚  Resources & Next Steps", font_size=40, color=WHITE, bold=True)
add_shape(slide, Inches(1), Inches(1.0), Inches(2), Inches(0.06), TEAL)

resources = [
    ("📄", "Download the code:", "github.com/saltcotraining-geoai/qgis-ai-site-selection"),
    ("📺", "Watch the demo video:", "youtube.com/@Saltco_GeoTraining"),
    ("📖", "QGIS PyQGIS Documentation:", "qgis.org/pyqgis"),
    ("📦", "Processing Algorithms Reference:", "docs.qgis.org/latest/en/docs/user_manual/processing_algs.html"),
    ("🐘", "PostGIS Documentation:", "postgis.net/documentation"),
    ("🌐", "GeoServer:", "geoserver.org"),
    ("📚", "Dataset:", "QGIS Training Manual (CC-BY-SA) + OpenStreetMap (ODbL)"),
]

for i, (emoji, label, link) in enumerate(resources):
    y = Inches(1.4) + i * Inches(0.85)
    add_shape(slide, Inches(1), y, Inches(0.08), Inches(0.6), TEAL)
    add_text_box(slide, Inches(1.3), y + Inches(0.02), Inches(3), Inches(0.4),
                 f"{emoji}  {label}", font_size=18, color=WHITE, bold=True)
    add_text_box(slide, Inches(4.5), y + Inches(0.02), Inches(7), Inches(0.4),
                 link, font_size=16, color=TEAL)

add_text_box(slide, Inches(1), Inches(6.2), Inches(11), Inches(0.5),
             "📺  Subscribe for more QGIS + AI automation workflows",
             font_size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

add_footer(slide)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 15: THANK YOU / Q&A
# ═══════════════════════════════════════════════════════════════════
slide = new_slide(15)
add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(2),
             "Thank You", font_size=72, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_shape(slide, Inches(5.5), Inches(3.5), Inches(2.3), Inches(0.06), TEAL)
add_text_box(slide, Inches(1), Inches(4), Inches(11), Inches(1),
             "Questions?", font_size=32, color=TEAL, align=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(5.2), Inches(11), Inches(0.5),
             '"The best GIS analyst is the one who automates their workflow."',
             font_size=18, color=GRAY, align=PP_ALIGN.CENTER, font_name="Calibri")

output_path = os.path.join(current_dir, "QGIS_AI_Site_Selection_Presentation.pptx")
prs.save(output_path)
print(f"✅ Presentation saved: {output_path}")
print(f"   {os.path.getsize(output_path) / 1024:.0f} KB — {len(prs.slides)} slides")
